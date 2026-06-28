import os
import re
import shutil
import sys
import tempfile
import zipfile

site = os.path.join(tempfile.gettempdir(), "codex_ppt_tools", "site")
sys.path.insert(0, site)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = os.path.abspath(os.path.join("output", "hangdun_risk_control_storyboard.pptx"))
os.makedirs(os.path.dirname(OUT), exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

NAVY = RGBColor(14, 45, 79)
BLUE = RGBColor(0, 91, 172)
CYAN = RGBColor(0, 158, 196)
TEAL = RGBColor(0, 128, 126)
ORANGE = RGBColor(232, 139, 39)
RED = RGBColor(199, 67, 58)
GREEN = RGBColor(72, 145, 96)
GRAY = RGBColor(92, 105, 117)
LIGHT = RGBColor(244, 248, 251)
WHITE = RGBColor(255, 255, 255)
PALE_BLUE = RGBColor(231, 241, 250)
PALE_TEAL = RGBColor(228, 245, 244)
PALE_ORANGE = RGBColor(253, 242, 229)
BORDER = RGBColor(207, 220, 232)
FONT = "Microsoft YaHei"


slides = [
    {
        "kind": "cover",
        "title": "航盾是客户经营全旅程线上风险防控的统一决策中枢",
        "subtitle": "航盾风控平台建设汇报",
        "note": "各位领导好，我是南航营销委航盾风控平台项目经理。本次主要汇报航盾是什么、为什么建、26年怎么建，以及需要哪些支持。",
    },
    {
        "kind": "agenda",
        "title": "本次汇报按“平台简介、风险清单、立项情况、资源需求”四部分展开",
        "items": ["航盾平台简介", "客营风险清单", "航盾26年立项情况", "资源需求与支持事项"],
        "note": "这一页先说明汇报路径，整体保持和原始材料一致，便于领导建立完整认知。",
    },
    {
        "kind": "what",
        "title": "航盾是面向客户经营全旅程的统一线上风险决策中枢，但不替代完整安全体系",
        "note": "首次汇报先讲清楚航盾定位。航盾是一道线上风控防线，承担识别、判断和决策能力，但不替代线下核验、内部审计、反诈宣传等完整安全体系。",
    },
    {
        "kind": "risk4",
        "title": "客户经营风险正在从零散个案变成跨场景、自动化、规模化风险",
        "cards": [
            ("黑产技术升级", "AI生成虚假身份、批量注册账号、自动化薅羊毛等攻击更隐蔽高效"),
            ("营销风险激增", "优惠券套现、虚假交易、恶意刷单造成营销费用无效流失"),
            ("组织流程痛点", "各业务团队独立风控，标准不一，形成盲区与重复建设"),
            ("核心能力短板", "依赖人工事后审核，发现周期长，处理效率和准确性受限"),
        ],
        "note": "这一页讲建设必要性。重点不是单个风险事件，而是风险形态已经变化，传统分散和人工处置难以支撑。",
    },
    {
        "kind": "value3",
        "title": "航盾通过全局数据、统一标准和复用能力，解决分散风控的效率与一致性问题",
        "cards": [
            ("全局数据沉淀", "形成跨业务风险视野，实现一处风险、全平台感知"),
            ("统一风险口径", "解决不同业务线标准不一、用户体验混乱的问题"),
            ("一次建设复用", "业务系统接入航盾API即可获得完整风控能力，降低重复建设成本"),
        ],
        "note": "这一页从管理视角讲航盾价值：减少重复建设，统一客户经营风险判断口径，并降低跨业务漏洞。",
    },
    {
        "kind": "model",
        "title": "航盾建设的核心是让业务系统专注采集与执行，让航盾统一判断与决策",
        "note": "这一页承接航盾定位，说明总体建设思路。业务系统负责接触用户、采集数据和执行处置；航盾集中沉淀数据、统一策略并输出风险决策。",
    },
    {
        "kind": "flow",
        "title": "业务系统负责采集和执行，航盾负责判断和决策，双方分工清晰",
        "steps": ["用户行为", "业务系统采集", "航盾风险判断", "返回处置指令", "业务系统执行"],
        "note": "这一页更落地解释业务接入后的协作方式。业务系统不需要自建完整风险判断能力，而是把关键数据传给航盾，并按照航盾指令执行。",
    },
    {
        "kind": "numbers",
        "title": "客营风险清单已收集109个需求，其中14条重点风险尚未实现自动化智能拦截",
        "note": "这一页说明航盾建设不是凭空提出，而是来自客营真实风险需求。目前清单中涉及账号安全、会员资产盗用、里程倒卖、钻漏洞获利等重点风险。",
    },
    {
        "kind": "journey",
        "title": "26年航盾将围绕关键用户旅程建立统一风险识别能力",
        "steps": ["账号登录", "受让人管理", "里程支付", "机票退改", "订单查询", "商城下单"],
        "note": "这是建设范围核心页。26年优先覆盖最容易产生资产和交易风险的关键环节，包括登录、里程支付、退改、查询和下单等场景。",
    },
    {
        "kind": "ops3",
        "title": "除场景接入外，航盾还将优化策略管理、可视化展示和风险处置效率",
        "cards": [
            ("策略配置", "优化规则配置流程，降低策略迭代门槛"),
            ("可视化展示", "完善风险态势、命中情况和处置结果展示"),
            ("风险处置", "提升处置协同效率，减少人为出错概率"),
        ],
        "note": "这一页补充说明项目不只是接接口，也包括提升后续风控运营和管理效率，保障风控业务高效、有序开展。",
    },
    {
        "kind": "gap",
        "title": "常驻地、IP地域和设备指纹等关键维度缺失，会影响风险识别准确性",
        "note": "这里先讲能力缺口，再自然引出预算和采购支持。常驻地、小程序端和H5端设备信息属于风控关键数据，缺失会影响后续风险识别准确性。",
    },
    {
        "kind": "support",
        "title": "为保障26年建设效果，需支持关键数据能力补齐、采购节奏协调和预算缺口解决",
        "note": "收尾需要明确支持事项：一是支持外采IP所在地域和设备指纹能力，二是支持协调采购窗口，三是支持解决约30万元预算缺口。以上请领导指导。",
    },
]


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color=BORDER, width=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_textbox(slide, text, x, y, w, h, font_size=18, color=NAVY, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_title(slide, text):
    add_textbox(slide, text, 0.65, 0.34, 11.8, 0.82, 23, NAVY, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.2), Inches(1.25), Inches(0.06))
    set_fill(line, ORANGE)
    line.line.fill.background()


def add_footer(slide, idx):
    add_textbox(slide, "南航营销委 | 航盾风控平台", 0.65, 7.05, 3.4, 0.25, 8, GRAY)
    add_textbox(slide, f"{idx:02d}", 12.25, 7.02, 0.45, 0.25, 9, GRAY, align=PP_ALIGN.RIGHT)


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.11))
    set_fill(band, BLUE)
    band.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.11), Inches(13.333), Inches(0.035))
    set_fill(accent, CYAN)
    accent.line.fill.background()


def add_card(slide, x, y, w, h, title, body, fill=LIGHT, border=BORDER, title_color=NAVY, body_color=GRAY, idx=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(box, fill)
    set_line(box, border, 1)
    tx = x + 0.24
    if idx is not None:
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.18), Inches(0.38), Inches(0.38))
        set_fill(circ, BLUE)
        circ.line.fill.background()
        add_textbox(slide, str(idx), x + 0.205, y + 0.225, 0.33, 0.22, 9, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        tx = x + 0.65
    add_textbox(slide, title, tx, y + 0.18, w - (tx - x) - 0.18, 0.33, 14, title_color, True)
    if body:
        add_textbox(slide, body, x + 0.24, y + 0.68, w - 0.48, h - 0.85, 10.5, body_color)
    return box


def add_connector(slide, x1, y1, x2, y2, color=BLUE):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(1.6)
    return c


notes = []
for idx, spec in enumerate(slides, start=1):
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    notes.append(spec["note"])
    kind = spec["kind"]

    if kind == "cover":
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        set_fill(rect, NAVY)
        rect.line.fill.background()
        for x1, y1, x2, y2 in [(8.4, 1.1, 11.8, 2.2), (7.6, 2.7, 12.3, 4.4), (8.9, 5.5, 12.0, 3.2), (6.8, 4.5, 10.0, 1.7)]:
            add_connector(slide, x1, y1, x2, y2, RGBColor(72, 154, 191))
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x2 - 0.06), Inches(y2 - 0.06), Inches(0.12), Inches(0.12))
            set_fill(dot, ORANGE)
            dot.line.fill.background()
        add_textbox(slide, spec["subtitle"], 0.85, 1.25, 6.2, 0.55, 22, RGBColor(214, 231, 244))
        add_textbox(slide, spec["title"], 0.82, 2.05, 8.0, 1.8, 34, WHITE, True)
        add_textbox(slide, "汇报人：南航营销委航盾风控平台项目经理", 0.88, 5.25, 6.3, 0.35, 15, RGBColor(219, 231, 240))
        add_textbox(slide, "2026年", 0.88, 5.72, 2.0, 0.3, 13, RGBColor(219, 231, 240))
        shield = slide.shapes.add_shape(MSO_SHAPE.HEPTAGON, Inches(9.15), Inches(2.35), Inches(2.0), Inches(2.15))
        set_fill(shield, RGBColor(18, 91, 136))
        shield.line.color.rgb = CYAN
        shield.line.width = Pt(2)
        add_textbox(slide, "航盾", 9.58, 3.13, 1.1, 0.42, 22, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        continue

    add_title(slide, spec["title"])
    add_footer(slide, idx)

    if kind == "agenda":
        xs = [0.85, 3.92, 6.99, 10.06]
        colors = [PALE_BLUE, PALE_TEAL, PALE_ORANGE, LIGHT]
        for i, item in enumerate(spec["items"]):
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(xs[i]), Inches(2.25), Inches(2.25), Inches(2.45))
            set_fill(box, colors[i])
            set_line(box, BORDER, 1.1)
            num = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(xs[i] + 0.75), Inches(1.72), Inches(0.75), Inches(0.75))
            set_fill(num, [BLUE, TEAL, ORANGE, NAVY][i])
            num.line.fill.background()
            add_textbox(slide, f"{i+1}", xs[i] + 0.98, 1.91, 0.3, 0.25, 15, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
            add_textbox(slide, item, xs[i] + 0.2, 3.0, 1.85, 0.55, 18, NAVY, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
            if i < 3:
                add_connector(slide, xs[i] + 2.33, 3.45, xs[i + 1] - 0.15, 3.45, BORDER)
        for x, label in zip([1.25, 4.25, 7.35, 10.4], ["认知建立", "需求来源", "建设计划", "支持事项"]):
            add_textbox(slide, label, x, 5.2, 1.6, 0.3, 11, GRAY, align=PP_ALIGN.CENTER)
    elif kind == "what":
        center = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.35), Inches(2.05), Inches(2.25), Inches(2.25))
        set_fill(center, BLUE)
        center.line.color.rgb = CYAN
        center.line.width = Pt(2)
        add_textbox(slide, "航盾", 5.82, 2.74, 1.3, 0.45, 25, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        for text, x, y in [("一体化", 1.1, 2.15), ("智能化", 1.1, 3.25), ("全流程", 1.1, 4.35), ("线上识别", 9.4, 2.0), ("风险判断", 9.4, 3.05), ("策略决策", 9.4, 4.1)]:
            add_card(slide, x, y, 2.2, 0.72, text, "", PALE_BLUE if x < 5 else PALE_TEAL, BORDER, BLUE)
            add_connector(slide, x + 2.2 if x < 5 else x, y + 0.36, 5.35 if x < 5 else 7.6, 3.18, BORDER)
        add_card(slide, 0.95, 5.45, 5.45, 0.9, "航盾覆盖", "南航用户旅程中的线上风险识别、防控与处置决策。", LIGHT, BORDER, BLUE)
        add_card(slide, 6.85, 5.45, 5.45, 0.9, "边界之外", "社工欺诈、线下冒用、钓鱼网站、内部人员与物理安全等需配套体系协同。", PALE_ORANGE, BORDER, ORANGE)
    elif kind == "risk4":
        coords = [(0.9, 1.75), (6.85, 1.75), (0.9, 4.15), (6.85, 4.15)]
        fills = [PALE_BLUE, PALE_ORANGE, PALE_TEAL, LIGHT]
        for i, (title, body) in enumerate(spec["cards"]):
            add_card(slide, coords[i][0], coords[i][1], 5.3, 1.55, title, body, fills[i], BORDER, [BLUE, ORANGE, TEAL, NAVY][i], idx=i + 1)
        mid = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.82), Inches(3.16), Inches(1.55), Inches(1.0))
        set_fill(mid, WHITE)
        set_line(mid, BORDER, 1.2)
        add_textbox(slide, "风险压力\n上升", 6.03, 3.36, 1.14, 0.38, 12, NAVY, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    elif kind == "value3":
        xs = [0.85, 4.65, 8.45]
        for i, (t, b) in enumerate(spec["cards"]):
            add_card(slide, xs[i], 2.05, 3.25, 2.35, t, b, [PALE_BLUE, PALE_TEAL, PALE_ORANGE][i], BORDER, [BLUE, TEAL, ORANGE][i], idx=i + 1)
        add_textbox(slide, "分散建设、口径不一、风险孤岛", 1.1, 5.35, 4.9, 0.35, 14, RED, True, PP_ALIGN.CENTER)
        add_connector(slide, 6.15, 5.5, 7.0, 5.5, BORDER)
        add_textbox(slide, "统一沉淀、统一判断、统一复用", 7.1, 5.35, 4.9, 0.35, 14, GREEN, True, PP_ALIGN.CENTER)
    elif kind == "model":
        add_card(slide, 0.9, 2.35, 3.0, 2.2, "业务系统", "负责用户触点、数据采集、前端提示与处置执行。", PALE_BLUE, BORDER, BLUE)
        add_card(slide, 5.15, 2.05, 3.0, 2.8, "航盾决策中枢", "统一沉淀风险数据，配置风险策略，输出识别结论与处置指令。", PALE_TEAL, BORDER, TEAL)
        add_card(slide, 9.4, 2.35, 3.0, 2.2, "处置结果", "拦截、二次验证、人工审核、放行等处置动作。", PALE_ORANGE, BORDER, ORANGE)
        add_connector(slide, 3.9, 3.45, 5.15, 3.45, BLUE)
        add_connector(slide, 8.15, 3.45, 9.4, 3.45, BLUE)
    elif kind == "flow":
        xs = [0.55, 3.0, 5.45, 7.9, 10.35]
        for i, st in enumerate(spec["steps"]):
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(xs[i]), Inches(2.75), Inches(1.9), Inches(1.15))
            set_fill(shape, [PALE_BLUE, PALE_BLUE, PALE_TEAL, PALE_TEAL, PALE_ORANGE][i])
            set_line(shape, BORDER, 1)
            add_textbox(slide, st, xs[i] + 0.12, 3.07, 1.66, 0.32, 13, [BLUE, BLUE, TEAL, TEAL, ORANGE][i], True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
            if i < 4:
                add_connector(slide, xs[i] + 1.9, 3.33, xs[i + 1], 3.33, BLUE)
        add_textbox(slide, "业务侧", 3.08, 4.35, 1.75, 0.28, 12, BLUE, True, PP_ALIGN.CENTER)
        add_textbox(slide, "航盾侧", 5.55, 4.35, 4.05, 0.28, 12, TEAL, True, PP_ALIGN.CENTER)
        add_textbox(slide, "业务侧", 10.42, 4.35, 1.75, 0.28, 12, ORANGE, True, PP_ALIGN.CENTER)
    elif kind == "numbers":
        add_textbox(slide, "109", 1.15, 2.05, 2.0, 0.85, 44, BLUE, True, PP_ALIGN.CENTER)
        add_textbox(slide, "个需求已收集", 1.08, 3.0, 2.2, 0.35, 14, GRAY, True, PP_ALIGN.CENTER)
        add_textbox(slide, "14", 4.0, 2.05, 1.6, 0.85, 44, ORANGE, True, PP_ALIGN.CENTER)
        add_textbox(slide, "条重点风险未自动化拦截", 3.5, 3.0, 2.7, 0.35, 14, GRAY, True, PP_ALIGN.CENTER)
        for i, (t, c) in enumerate([("账号安全", BLUE), ("会员资产盗用", TEAL), ("里程倒卖", ORANGE), ("钻漏洞获利", RED)]):
            add_card(slide, 7.0, 1.75 + i * 0.92, 4.65, 0.68, t, "", LIGHT, BORDER, c)
        add_textbox(slide, "重点对接场景：登录、里程支付、下单", 1.05, 4.8, 10.9, 0.45, 16, NAVY, True, PP_ALIGN.CENTER)
        add_textbox(slide, "数据来源：客营风险清单，在线清单链接需正式确认", 1.05, 5.45, 10.9, 0.3, 10, GRAY, False, PP_ALIGN.CENTER)
    elif kind == "journey":
        y = 3.35
        add_connector(slide, 1.15, y, 12.0, y, BORDER)
        xs = [1.15, 3.25, 5.35, 7.45, 9.55, 11.65]
        for i, st in enumerate(spec["steps"]):
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(xs[i] - 0.22), Inches(y - 0.22), Inches(0.44), Inches(0.44))
            set_fill(dot, [BLUE, TEAL, ORANGE, BLUE, TEAL, ORANGE][i])
            dot.line.fill.background()
            add_textbox(slide, str(i + 1), xs[i] - 0.09, y - 0.1, 0.18, 0.13, 8, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
            add_textbox(slide, st, xs[i] - 0.75, y + 0.45, 1.5, 0.5, 11.5, NAVY, True, PP_ALIGN.CENTER)
            add_textbox(slide, "风险识别", xs[i] - 0.55, y - 1.05, 1.1, 0.28, 9, GRAY, False, PP_ALIGN.CENTER)
        add_textbox(slide, "围绕用户旅程关键节点，逐步形成统一风险识别能力", 2.1, 5.35, 9.1, 0.35, 15, BLUE, True, PP_ALIGN.CENTER)
    elif kind == "ops3":
        xs = [0.95, 4.75, 8.55]
        for i, (t, b) in enumerate(spec["cards"]):
            add_card(slide, xs[i], 2.12, 3.05, 2.55, t, b, [PALE_BLUE, PALE_TEAL, PALE_ORANGE][i], BORDER, [BLUE, TEAL, ORANGE][i], idx=i + 1)
        add_textbox(slide, "降低操作门槛  |  减少人为错误  |  提升迭代效率", 2.05, 5.45, 9.1, 0.42, 16, NAVY, True, PP_ALIGN.CENTER)
    elif kind == "gap":
        add_card(slide, 0.95, 2.05, 3.1, 2.0, "当前采集不足", "业务系统无法完整采集用户常驻地，以及小程序端、H5端设备信息。", PALE_BLUE, BORDER, BLUE)
        add_card(slide, 5.15, 2.05, 3.1, 2.0, "判断维度缺失", "关键数据缺口会影响异常登录、资产交易和下单等场景的识别准确性。", PALE_TEAL, BORDER, TEAL)
        add_card(slide, 9.35, 2.05, 3.1, 2.0, "需补齐能力", "外采IP所在地域和设备指纹数据产品，补全风险识别基础维度。", PALE_ORANGE, BORDER, ORANGE)
        add_connector(slide, 4.05, 3.05, 5.15, 3.05, BLUE)
        add_connector(slide, 8.25, 3.05, 9.35, 3.05, BLUE)
        for i, t in enumerate(["用户常驻地", "IP所在地域", "设备指纹", "小程序/H5设备信息"]):
            add_textbox(slide, t, 1.45 + i * 2.72, 5.25, 2.1, 0.35, 12, NAVY, True, PP_ALIGN.CENTER)
    elif kind == "support":
        items = [
            ("支持外采能力", "补齐IP所在地域与设备指纹两类关键数据能力"),
            ("支持采购协调", "关注采购窗口，避免影响26年建设节奏"),
            ("支持预算缺口", "支持解决约30万元费用缺口，保障能力落地"),
        ]
        xs = [0.9, 4.75, 8.6]
        for i, (t, b) in enumerate(items):
            add_card(slide, xs[i], 1.85, 3.05, 2.25, t, b, [PALE_BLUE, PALE_TEAL, PALE_ORANGE][i], BORDER, [BLUE, TEAL, ORANGE][i], idx=i + 1)
        add_textbox(slide, "下一步：完善测算依据、细化场景接入排期、推进重点场景对接", 1.15, 5.15, 11.0, 0.45, 16, NAVY, True, PP_ALIGN.CENTER)
        add_textbox(slide, "注：预算总额“432万”表述需正式确认后用于定稿", 1.15, 5.78, 11.0, 0.3, 10, GRAY, False, PP_ALIGN.CENTER)

prs.save(OUT)


def add_notes_to_pptx(path, slide_notes):
    temp_dir = os.path.join(tempfile.gettempdir(), "ppt_notes_pkg")
    if os.path.exists(temp_dir):
        shutil.rmtree(temp_dir)
    os.makedirs(temp_dir)
    with zipfile.ZipFile(path, "r") as z:
        z.extractall(temp_dir)

    ppt_dir = os.path.join(temp_dir, "ppt")
    notes_dir = os.path.join(ppt_dir, "notesSlides")
    rels_dir = os.path.join(notes_dir, "_rels")
    os.makedirs(notes_dir, exist_ok=True)
    os.makedirs(rels_dir, exist_ok=True)

    ct_path = os.path.join(temp_dir, "[Content_Types].xml")
    with open(ct_path, "r", encoding="utf-8") as f:
        ct = f.read()
    if "notesSlide+xml" not in ct:
        override = "".join(
            f'<Override PartName="/ppt/notesSlides/notesSlide{i}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>'
            for i in range(1, len(slide_notes) + 1)
        )
        ct = ct.replace("</Types>", override + "</Types>")
        with open(ct_path, "w", encoding="utf-8") as f:
            f.write(ct)

    nm_dir = os.path.join(ppt_dir, "notesMasters")
    nm_rels = os.path.join(nm_dir, "_rels")
    os.makedirs(nm_dir, exist_ok=True)
    os.makedirs(nm_rels, exist_ok=True)
    nm_path = os.path.join(nm_dir, "notesMaster1.xml")
    if not os.path.exists(nm_path):
        with open(nm_path, "w", encoding="utf-8") as f:
            f.write(
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><p:notesMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr></p:spTree></p:cSld><p:clrMap accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" bg1="lt1" bg2="lt2" folHlink="folHlink" hlink="hlink" tx1="dk1" tx2="dk2"/></p:notesMaster>'
            )
        with open(os.path.join(nm_rels, "notesMaster1.xml.rels"), "w", encoding="utf-8") as f:
            f.write('<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>')

    pres_rels_path = os.path.join(ppt_dir, "_rels", "presentation.xml.rels")
    with open(pres_rels_path, "r", encoding="utf-8") as f:
        pres_rels = f.read()
    if "notesMaster" not in pres_rels:
        ids = [int(x) for x in re.findall(r'Id="rId(\d+)"', pres_rels)] or [0]
        rid = max(ids) + 1
        pres_rels = pres_rels.replace(
            "</Relationships>",
            f'<Relationship Id="rId{rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster" Target="notesMasters/notesMaster1.xml"/></Relationships>',
        )
        with open(pres_rels_path, "w", encoding="utf-8") as f:
            f.write(pres_rels)

    def esc(s):
        return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")

    for i, note in enumerate(slide_notes, start=1):
        slide_rels_path = os.path.join(ppt_dir, "slides", "_rels", f"slide{i}.xml.rels")
        with open(slide_rels_path, "r", encoding="utf-8") as f:
            rels = f.read()
        ids = [int(x) for x in re.findall(r'Id="rId(\d+)"', rels)] or [0]
        rid = max(ids) + 1
        if f"notesSlides/notesSlide{i}.xml" not in rels:
            rels = rels.replace(
                "</Relationships>",
                f'<Relationship Id="rId{rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide" Target="../notesSlides/notesSlide{i}.xml"/></Relationships>',
            )
            with open(slide_rels_path, "w", encoding="utf-8") as f:
                f.write(rels)

        note_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld><p:spTree>
    <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
    <p:sp><p:nvSpPr><p:cNvPr id="2" name="Slide Image Placeholder 1"/><p:cNvSpPr><a:spLocks noGrp="1" noRot="1" noChangeAspect="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldImg" idx="1"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>
    <p:sp><p:nvSpPr><p:cNvPr id="3" name="Notes Placeholder 2"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="body" sz="quarter" idx="2"/></p:nvPr></p:nvSpPr><p:spPr><a:xfrm><a:off x="685800" y="4400000"/><a:ext cx="7772400" cy="2057400"/></a:xfrm></p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr lang="zh-CN" sz="1200"><a:latin typeface="Microsoft YaHei"/><a:ea typeface="Microsoft YaHei"/></a:rPr><a:t>{esc(note)}</a:t></a:r></a:p></p:txBody></p:sp>
  </p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:notes>'''
        with open(os.path.join(notes_dir, f"notesSlide{i}.xml"), "w", encoding="utf-8") as f:
            f.write(note_xml)
        note_rels = f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="../slides/slide{i}.xml"/><Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster" Target="../notesMasters/notesMaster1.xml"/></Relationships>'
        with open(os.path.join(rels_dir, f"notesSlide{i}.xml.rels"), "w", encoding="utf-8") as f:
            f.write(note_rels)

    tmp_path = path + ".tmp"
    with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as z:
        for root, _, files in os.walk(temp_dir):
            for file in files:
                full = os.path.join(root, file)
                arc = os.path.relpath(full, temp_dir).replace(os.sep, "/")
                z.write(full, arc)
    shutil.move(tmp_path, path)
    shutil.rmtree(temp_dir)


add_notes_to_pptx(OUT, notes)
print(OUT)
